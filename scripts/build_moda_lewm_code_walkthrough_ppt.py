from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_PPT = ROOT / "docs" / "moda_lewm_code_walkthrough_report.pptx"
OUT_SCRIPT = ROOT / "docs" / "moda_lewm_code_walkthrough_report_script.md"


WIDE = (13.333, 7.5)

COLORS = {
    "navy": RGBColor(22, 36, 57),
    "blue": RGBColor(31, 91, 180),
    "cyan": RGBColor(46, 139, 192),
    "green": RGBColor(46, 125, 87),
    "orange": RGBColor(203, 121, 43),
    "red": RGBColor(181, 68, 68),
    "gray": RGBColor(92, 102, 112),
    "light": RGBColor(245, 247, 250),
    "line": RGBColor(215, 221, 230),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(28, 32, 36),
}


def set_text(frame, text, font_size=20, bold=False, color="black", align=None):
    frame.clear()
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "PingFang SC"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return p


def add_textbox(slide, x, y, w, h, text, font_size=20, bold=False, color="black", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.08)
    box.text_frame.margin_right = Inches(0.08)
    box.text_frame.margin_top = Inches(0.04)
    box.text_frame.margin_bottom = Inches(0.04)
    set_text(box.text_frame, text, font_size, bold, color, align)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.35, 11.8, 0.55, title, 25, True, "navy")
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.62),
        Inches(0.95),
        Inches(1.2),
        Inches(0.055),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLORS["blue"]
    slide.shapes[-1].line.fill.background()
    if subtitle:
        add_textbox(slide, 0.55, 1.05, 11.8, 0.36, subtitle, 11, False, "gray")


def add_footer(slide, idx):
    add_textbox(slide, 0.55, 7.13, 7.0, 0.22, "LeWM + MoDA technical path", 8, False, "gray")
    add_textbox(slide, 12.25, 7.13, 0.5, 0.22, str(idx), 8, False, "gray", PP_ALIGN.RIGHT)


def add_chip(slide, x, y, w, h, text, fill="light", text_color="navy", font_size=13):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    set_text(tf, text, font_size, True, text_color, PP_ALIGN.CENTER)
    return shape


def add_bullets(slide, x, y, w, h, bullets, font_size=16, color="black"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "PingFang SC"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(7)
    return box


def add_card(slide, x, y, w, h, title, body, accent="blue"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.08),
        Inches(h),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLORS[accent]
    slide.shapes[-1].line.fill.background()
    add_textbox(slide, x + 0.22, y + 0.16, w - 0.35, 0.32, title, 14, True, "navy")
    add_textbox(slide, x + 0.22, y + 0.55, w - 0.35, h - 0.65, body, 11, False, "black")
    return shape


def add_arrow(slide, x1, y1, x2, y2, color="gray"):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(1.8)
    return line


def add_table(slide, x, y, w, h, headers, rows, font_size=10):
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    ).table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["navy"]
        set_text(cell.text_frame, header, font_size, True, "white", PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["light"] if i % 2 else COLORS["white"]
            set_text(cell.text_frame, str(val), font_size, False, "black")
    return table


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(WIDE[0])
    prs.slide_height = Inches(WIDE[1])
    blank = prs.slide_layouts[6]
    slides = []

    # 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COLORS["light"]
    add_textbox(s, 0.75, 0.75, 9.8, 0.85, "LeWM + MoDA 技术路径汇报", 34, True, "navy")
    add_textbox(s, 0.78, 1.68, 10.8, 0.45, "从结构接入到 planning 诊断，再到 MoDA-only residual proposal", 18, False, "gray")
    add_chip(s, 0.78, 2.55, 2.2, 0.42, "LeWM 架构", "white")
    add_chip(s, 3.25, 2.55, 2.2, 0.42, "MoDA 接入", "white")
    add_chip(s, 5.72, 2.55, 2.45, 0.42, "PushT Planning", "white")
    add_chip(s, 8.42, 2.55, 2.7, 0.42, "Residual Proposal", "white")
    add_textbox(s, 0.8, 5.9, 10.5, 0.5, "核心结论：MoDA 有候选覆盖价值，但 raw cost 与 planning success 不完全对齐；当前最可信的 MoDA-only 改进方向是 action proposal correction。", 17, True, "navy")
    slides.append(s)

    # 2
    s = prs.slides.add_slide(blank)
    add_title(s, "任务背景：LeWM 不是单任务框架，但当前实验主线聚焦 PushT")
    rows = [
        ["PushT", "推 T 形物体到目标位姿", "主实验任务：CEM、candidate pool、near-miss failure 都围绕它"],
        ["Cube", "方块/物体运动控制", "后续可用于物体操控泛化验证"],
        ["TwoRooms", "跨房间/空间导航", "动力学和 action geometry 与 PushT 不同"],
        ["Reacher", "连续控制 reaching", "经典控制任务，可作为轻量泛化检查"],
    ]
    add_table(s, 0.75, 1.45, 11.8, 2.35, ["任务", "大概含义", "在本项目中的位置"], rows, 10)
    add_card(
        s,
        0.85,
        4.25,
        5.6,
        1.35,
        "汇报边界",
        "LeWM 官方/相关 checkpoint 可覆盖 PushT、Cube、TwoRooms、Reacher；但当前 MoDA 接入、诊断和结果主要在 PushT 上完成。",
        "orange",
    )
    add_card(
        s,
        6.85,
        4.25,
        5.2,
        1.35,
        "为什么先做 PushT",
        "PushT 有连续动作、接触动力学和明显 near-miss failure，适合分析 candidate coverage、cost-success mismatch 和 proposal correction。",
        "blue",
    )
    add_footer(s, 2)
    slides.append(s)

    # 3
    s = prs.slides.add_slide(blank)
    add_title(s, "LeWM 基本架构：Latent World Model + CEM Planning")
    add_chip(s, 0.8, 1.55, 2.1, 0.55, "Observation", "light")
    add_arrow(s, 2.95, 1.83, 3.65, 1.83, "gray")
    add_chip(s, 3.7, 1.55, 2.1, 0.55, "Encoder", "light")
    add_arrow(s, 5.85, 1.83, 6.55, 1.83, "gray")
    add_chip(s, 6.6, 1.55, 2.35, 0.55, "Latent State", "light")
    add_arrow(s, 8.98, 1.83, 9.68, 1.83, "gray")
    add_chip(s, 9.75, 1.55, 2.35, 0.55, "Predictor", "light")
    add_card(s, 0.8, 2.75, 3.55, 1.55, "World Model", "在 latent 空间预测未来状态，不直接在 pixel/state raw space 里做长 rollout。", "blue")
    add_card(s, 4.72, 2.75, 3.55, 1.55, "CEM Planner", "采样 action sequence，rollout 后计算 cost，保留 elites 并更新采样分布。", "green")
    add_card(s, 8.65, 2.75, 3.55, 1.55, "Planning Output", "最终选 cost 最低的 action candidate；这就是 top1 success 的关键位置。", "orange")
    add_bullets(
        s,
        0.95,
        5.0,
        11.3,
        1.1,
        [
            "LeWM 的优势：把高维观测压到 latent dynamics，再用 CEM 做模型预测控制。",
            "本项目关心的问题：MoDA 改 predictor/rollout 后，candidate pool 和 final top1 selection 是否真的改善。",
        ],
        15,
    )
    add_footer(s, 3)
    slides.append(s)

    # 4
    s = prs.slides.add_slide(blank)
    add_title(s, "MoDA 的核心优势：跨 depth retrieval 的 attention")
    add_card(s, 0.8, 1.35, 3.6, 1.55, "普通 Transformer Block", "每层主要基于当前层表示做 self-attention，depth 间信息通过逐层传递。", "gray")
    add_card(s, 4.85, 1.35, 3.6, 1.55, "MoDA Attention", "当前层 query 可以访问前面层累计的 K/V depth cache，形成 mixture-of-depth retrieval。", "blue")
    add_card(s, 8.9, 1.35, 3.35, 1.55, "对 LeWM 的意义", "预测未来 latent 时能检索跨层信息，理论上增强 state rollout 表达。", "green")
    add_chip(s, 1.15, 3.72, 1.45, 0.45, "Layer 1 K/V", "white")
    add_chip(s, 2.9, 3.72, 1.45, 0.45, "Layer 2 K/V", "white")
    add_chip(s, 4.65, 3.72, 1.45, 0.45, "Layer 3 K/V", "white")
    add_arrow(s, 6.25, 3.95, 7.15, 3.95, "gray")
    add_chip(s, 7.25, 3.62, 2.05, 0.65, "Depth Cache", "light")
    add_arrow(s, 9.42, 3.95, 10.2, 3.95, "gray")
    add_chip(s, 10.3, 3.62, 1.75, 0.65, "MoDA Kernel", "light")
    add_bullets(
        s,
        1.0,
        5.0,
        11.0,
        0.95,
        [
            "接入目标不是先改 planner，而是先让 LeWM predictor 具备 MoDA 的 depth-aware representation。",
            "后续再检验这种 representation 是否转化成更好的 planning candidate 和 top1 selection。",
        ],
        15,
    )
    add_footer(s, 4)
    slides.append(s)

    # 5
    s = prs.slides.add_slide(blank)
    add_title(s, "第一阶段：把 MoDA 接到 LeWM predictor，而不是重写整个系统")
    add_card(s, 0.75, 1.3, 3.65, 1.4, "`moda_module_exact.py`", "核心结构实现：保留 LeWM conditional block，只把 attention kernel 替换成 MoDA exact/Triton 路径。", "blue")
    add_card(s, 4.85, 1.3, 3.65, 1.4, "`MoDAAttentionExact`", "无 cache 时退回 causal attention；有 cached K/V 时调用 MoDA parallel kernel。", "green")
    add_card(s, 8.95, 1.3, 3.65, 1.4, "`MoDATransformerExact`", "逐层收集 K/V，构建 depth cache；用 cache_window/detach_cache 控制显存和稳定性。", "orange")
    add_bullets(
        s,
        0.95,
        3.35,
        11.3,
        1.6,
        [
            "保留 LeWM 原有 AdaLN-zero conditioning：shift/scale/gate 仍然调制 attention 和 MLP 分支。",
            "兼容不同运行环境：本地、5090 服务器、实习服务器都可以通过 MoDA Triton root 找 kernel。",
            "这一步的贡献是结构接入：让 LeWM predictor 能使用 MoDA depth cache，而不是直接宣称 planner 已经变强。",
        ],
        15,
    )
    add_footer(s, 5)
    slides.append(s)

    # 6
    s = prs.slides.add_slide(blank)
    add_title(s, "第二阶段：训练入口把 MoDA checkpoint 接入 LeWM pipeline")
    add_card(s, 0.75, 1.35, 3.55, 1.35, "`train_encoder_moda.py`", "读取 LeWM config，构造数据，初始化 MoDA encoder/predictor，保存 checkpoint。", "blue")
    add_card(s, 4.85, 1.35, 3.55, 1.35, "`train_moda.py`", "早期 prototype：验证 MoDA module 在 LeWM 框架下能训练、能出 checkpoint。", "green")
    add_card(s, 8.95, 1.35, 3.55, 1.35, "`train_moda_exact.py`", "exact kernel 版本：配合 `moda_module_exact.py`，更贴近 official MoDA kernel 路径。", "orange")
    add_chip(s, 1.05, 3.55, 1.95, 0.5, "Config", "light")
    add_arrow(s, 3.08, 3.8, 3.78, 3.8)
    add_chip(s, 3.9, 3.55, 2.15, 0.5, "Dataset", "light")
    add_arrow(s, 6.12, 3.8, 6.82, 3.8)
    add_chip(s, 6.95, 3.55, 2.2, 0.5, "MoDA Model", "light")
    add_arrow(s, 9.22, 3.8, 9.92, 3.8)
    add_chip(s, 10.05, 3.55, 2.0, 0.5, "Checkpoint", "light")
    add_bullets(
        s,
        0.95,
        4.95,
        11.3,
        0.95,
        [
            "训练入口不是最终方法本身，但它们让后续 stateroll candidate pool、PAC-MoDA 和 residual proposal 有可评估的 MoDA checkpoint。",
            "汇报时重点讲 pipeline 位置，不需要逐个训练参数展开。",
        ],
        15,
    )
    add_footer(s, 6)
    slides.append(s)

    # 7
    s = prs.slides.add_slide(blank)
    add_title(s, "Planning 评估问题：candidate pool 有成功候选，但 raw cost 未必选对")
    add_chip(s, 0.85, 1.45, 2.2, 0.5, "30 candidates / episode", "light")
    add_arrow(s, 3.15, 1.7, 3.95, 1.7)
    add_chip(s, 4.05, 1.45, 2.4, 0.5, "raw MoDA cost ranking", "light")
    add_arrow(s, 6.55, 1.7, 7.35, 1.7)
    add_chip(s, 7.45, 1.45, 2.15, 0.5, "top1 action", "light")
    add_arrow(s, 9.72, 1.7, 10.45, 1.7)
    add_chip(s, 10.55, 1.45, 1.7, 0.5, "success?", "light")
    add_card(s, 0.9, 2.65, 3.6, 1.55, "Oracle 不低", "同一个 episode 的候选池里经常存在成功 candidate，说明 MoDA 并不是完全没有覆盖到好动作。", "green")
    add_card(s, 4.85, 2.65, 3.6, 1.55, "Top1 不强", "raw cost 往往把 near-miss failure 排到第一，导致最终 action 失败。", "red")
    add_card(s, 8.8, 2.65, 3.6, 1.55, "核心矛盾", "问题不是 only coverage，而是 cost-success alignment 和 candidate selection。", "orange")
    add_textbox(s, 1.0, 5.15, 11.0, 0.45, "所以后续实验全部围绕一个问题：怎么让 MoDA 在同一个 episode 内把成功候选选到 top1？", 18, True, "navy")
    add_footer(s, 7)
    slides.append(s)

    # 8
    s = prs.slides.add_slide(blank)
    add_title(s, "诊断一：为什么 BCE/global AUC 高，但 top1 不涨")
    rows = [
        ["episode-only", "global AUC 0.692", "top1 55.33", "intra-episode AUC 0.500"],
        ["candidate-only", "global AUC 0.703", "top1 55.50", "intra-episode AUC 0.559"],
        ["raw cost", "global AUC 0.494", "top1 55.33", "success-over-rank0 0.0%"],
    ]
    add_table(s, 0.85, 1.35, 11.5, 1.75, ["设置", "AUC", "Top1", "关键诊断"], rows, 10)
    add_card(
        s,
        0.9,
        3.65,
        5.45,
        1.45,
        "解释",
        "高 global AUC 很多来自 episode difficulty leakage：模型知道这个 episode 容不容易成功，但不知道同 episode 内哪个 candidate 更该排第一。",
        "orange",
    )
    add_card(
        s,
        6.85,
        3.65,
        5.2,
        1.45,
        "结论",
        "post-hoc feature-level reranking 很难稳定救 MoDA-only top1，必须进入 planner/candidate generation 层面。",
        "red",
    )
    add_footer(s, 8)
    slides.append(s)

    # 9
    s = prs.slides.add_slide(blank)
    add_title(s, "诊断二：我们系统排除了几条看似合理但不稳定的路线")
    rows = [
        ["Final rerank / BCE / selector", "AUC 可提升", "top1 基本不涨；intra-episode 不可分"],
        ["Planner-in-loop calibrated CEM", "J_plan = J_raw - lambda U", "小幅信号，但不稳定"],
        ["Action-sensitive contrastive", "拉开 success/near-miss embedding", "proxy 改善不等于 top1 改善"],
        ["Search scaling", "增加候选/步数/预算", "不能稳定修复 raw cost selection"],
    ]
    add_table(s, 0.75, 1.35, 11.9, 3.05, ["路线", "尝试内容", "结论"], rows, 10)
    add_textbox(
        s,
        0.9,
        5.05,
        11.4,
        0.65,
        "这个阶段的意义不是“都失败了”，而是定位 bottleneck：MoDA 的候选生成和 action proposal 需要被修正，不能只在最后 30 个候选上重排。",
        16,
        True,
        "navy",
    )
    add_footer(s, 9)
    slides.append(s)

    # 10
    s = prs.slides.add_slide(blank)
    add_title(s, "Baseline-safe / PAC-MoDA：能证明辅助价值，但不是 MoDA-only")
    add_card(s, 0.85, 1.35, 3.75, 1.55, "系统逻辑", "baseline planner 保底；MoDA 只在 gate 判断为低风险/有机会时介入。", "blue")
    add_card(s, 4.95, 1.35, 3.75, 1.55, "正向结果", "bsl-integrated 曾有约 81 -> 82/83 的系统级提升，说明 MoDA 候选有互补信息。", "green")
    add_card(s, 9.05, 1.35, 3.25, 1.55, "边界", "这依赖 baseline fallback，不能当成 MoDA standalone planner。", "orange")
    add_chip(s, 1.25, 3.75, 2.0, 0.5, "Baseline action", "light")
    add_arrow(s, 3.35, 4.0, 4.1, 4.0)
    add_chip(s, 4.2, 3.75, 1.6, 0.5, "Gate", "light")
    add_arrow(s, 5.9, 4.0, 6.65, 4.0)
    add_chip(s, 6.75, 3.75, 2.05, 0.5, "MoDA candidate", "light")
    add_arrow(s, 8.9, 4.0, 9.65, 4.0)
    add_chip(s, 9.75, 3.75, 1.9, 0.5, "Final action", "light")
    add_textbox(s, 0.95, 5.35, 11.1, 0.42, "汇报口径：这条线证明 MoDA auxiliary candidate source 有价值，但当前主线仍要区分 MoDA-only 与 baseline-safe integration。", 15, True, "navy")
    add_footer(s, 10)
    slides.append(s)

    # 11
    s = prs.slides.add_slide(blank)
    add_title(s, "当前最重要的 MoDA-only 正向方向：Success-conditioned Residual Proposal")
    add_textbox(s, 0.9, 1.25, 11.2, 0.55, "核心思想：不再只改 final ranking，而是学习把 raw rank0 near-miss failure 的 action 往成功 action 方向推。", 18, True, "navy")
    add_chip(s, 0.95, 2.25, 2.35, 0.55, "raw rank0 failure", "light")
    add_arrow(s, 3.42, 2.53, 4.12, 2.53)
    add_chip(s, 4.25, 2.25, 2.35, 0.55, "learn delta_a", "light")
    add_arrow(s, 6.72, 2.53, 7.42, 2.53)
    add_chip(s, 7.55, 2.25, 2.35, 0.55, "shifted candidates", "light")
    add_arrow(s, 10.02, 2.53, 10.72, 2.53)
    add_chip(s, 10.85, 2.25, 1.35, 0.55, "re-score", "light")
    add_card(
        s,
        1.0,
        3.55,
        5.25,
        1.45,
        "训练信号",
        "在 raw rank0 失败但同 episode 有成功候选的样本中，构造 delta_a = success_action - raw_rank0_action。",
        "green",
    )
    add_card(
        s,
        6.8,
        3.55,
        5.25,
        1.45,
        "评估方式",
        "在线 CEM 生成 raw candidates；预测 residual；生成 shifted candidates；再用 raw/calibrated cost 重新评估。",
        "blue",
    )
    add_footer(s, 11)
    slides.append(s)

    # 12
    s = prs.slides.add_slide(blank)
    add_title(s, "Residual Proposal 代码路径")
    add_card(s, 0.75, 1.25, 3.65, 1.35, "`train_data`", "从 stateroll/MoDA candidate pool 中抽取 raw rank0 failure -> best success 的 residual target。", "green")
    add_card(s, 4.85, 1.25, 3.65, 1.35, "`feature_np` / `online_feature`", "离线和在线构造一致的 cost、rank、trajectory、latent、action statistics 特征。", "blue")
    add_card(s, 8.95, 1.25, 3.65, 1.35, "`fit_ridge`", "用 ridge regression 学 residual；稳定、快速、可解释，适合验证 proposal correction 信号。", "orange")
    add_card(s, 0.75, 3.25, 3.65, 1.35, "`score_candidates`", "对 raw/shifted actions 重新 rollout，得到 raw cost、utility 和 calibrated plan score。", "blue")
    add_card(s, 4.85, 3.25, 3.65, 1.35, "`run_split`", "训练 residual model，生成 shifted candidates，比较 top1/top3/top5/oracle/near-miss。", "green")
    add_card(s, 8.95, 3.25, 3.65, 1.35, "`confirm50_audit`", "固定 eval indices，做 first20/added30/all50 和 scale sensitivity，避免小样本误读。", "red")
    add_footer(s, 12)
    slides.append(s)

    # 13
    s = prs.slides.add_slide(blank)
    add_title(s, "Residual Proposal 结果：方向正，但不能过度宣称")
    rows = [
        ["medium20", "raw 60", "residual 65", "+5", "near-miss 6 -> 3"],
        ["confirm50", "raw 46", "residual 49/50", "+3/+4", "near-miss 14 -> 10/11"],
    ]
    add_table(s, 1.0, 1.45, 11.1, 1.5, ["评估", "Raw MoDA", "Residual", "Paired gain", "Near-miss"], rows, 11)
    add_card(
        s,
        1.0,
        3.55,
        5.2,
        1.45,
        "可以说",
        "residual proposal 是目前唯一真正 MoDA-only 的正向方向，显示 consistent paired gain 和 near-miss reduction。",
        "green",
    )
    add_card(
        s,
        6.75,
        3.55,
        5.2,
        1.45,
        "不能说",
        "不能说 MoDA-only 已经稳定 65+；medium20 的 65 受 eval set 难度影响，confirm50 后更稳结论是小幅正向。",
        "orange",
    )
    add_footer(s, 13)
    slides.append(s)

    # 14
    s = prs.slides.add_slide(blank)
    add_title(s, "完整技术路径：从接入到诊断，再到 proposal correction")
    steps = [
        ("1. MoDA 结构接入", "`moda_module.py` / `moda_module_exact.py`"),
        ("2. 训练入口跑通", "`train_encoder_moda.py` / `train_moda_exact.py`"),
        ("3. Planning candidate pool", "stateroll / raw_rollout_npz / CEM candidates"),
        ("4. MoDA-only 诊断", "AUC leakage / intra-episode AUC / near-miss"),
        ("5. Baseline-safe 结果", "证明辅助价值，但不是 MoDA-only"),
        ("6. Residual proposal", "candidate generation correction"),
    ]
    y = 1.35
    for i, (t, b) in enumerate(steps):
        add_card(s, 0.85 + (i % 2) * 6.0, y + (i // 2) * 1.6, 5.35, 1.05, t, b, ["blue", "green", "orange", "red", "gray", "blue"][i])
    add_footer(s, 14)
    slides.append(s)

    # 15
    s = prs.slides.add_slide(blank)
    add_title(s, "汇报主结论")
    add_bullets(
        s,
        0.95,
        1.35,
        11.4,
        3.2,
        [
            "LeWM 本身是多任务 latent world model + CEM planning 框架，当前 MoDA 实验主要在 PushT 上做深。",
            "MoDA 的优势是 depth-aware representation / state rollout，但这个优势不会自动变成 planning top1 success。",
            "系统诊断说明：MoDA candidate pool 有成功候选，主要瓶颈是 raw cost 和 success 不对齐，以及 near-miss failure 抢 rank0。",
            "post-hoc rerank / AUC calibration 不能稳定解决同 episode 内 top1 selection。",
            "当前最可信的 MoDA-only 方向是 success-conditioned residual proposal：从 action proposal generation 层面修正候选。",
        ],
        17,
    )
    add_textbox(s, 1.0, 5.35, 11.0, 0.5, "一句话：MoDA 值得接入 LeWM，但要作为 planning-aware candidate generation correction，而不是只做最后重排。", 18, True, "navy")
    add_footer(s, 15)
    slides.append(s)

    # 16
    s = prs.slides.add_slide(blank)
    add_title(s, "后续工作：把当前正信号变成稳健方法")
    add_card(s, 0.9, 1.35, 3.65, 1.55, "Residual gate", "判断什么时候适合使用 residual correction，避免 raw success 被 shifted action harm。", "green")
    add_card(s, 4.85, 1.35, 3.65, 1.55, "Scale selection", "在同一批 eval indices 上选择 residual scale，减少 action shift 过大风险。", "blue")
    add_card(s, 8.8, 1.35, 3.65, 1.55, "Cross-task validation", "将同样诊断流程迁移到 Cube / TwoRooms / Reacher，检验泛化。", "orange")
    add_card(s, 0.9, 3.55, 3.65, 1.55, "Training-level objective", "未来可以把 near-miss rank loss / action-sensitive contrastive 放进 world model 训练。", "red")
    add_card(s, 4.85, 3.55, 3.65, 1.55, "Cleaner benchmark", "固定 indices、multi-seed、paired metrics，避免 absolute top1 被 eval set 难度误导。", "gray")
    add_card(s, 8.8, 3.55, 3.65, 1.55, "Release package", "代码、关键权重、candidate pool 和 walkthrough 已整理到 GitHub/LFS，便于复现和交接。", "blue")
    add_footer(s, 16)
    slides.append(s)

    # 17
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COLORS["navy"]
    add_textbox(s, 0.85, 1.0, 11.5, 0.8, "Takeaway", 34, True, "white")
    add_textbox(
        s,
        0.9,
        2.1,
        11.2,
        1.6,
        "MoDA 的问题不是没有候选价值，而是 raw cost 与 planning success 不对齐；\n我们从结构接入、训练、candidate pool 诊断一路定位到 candidate generation bottleneck。",
        22,
        True,
        "white",
    )
    add_textbox(
        s,
        0.9,
        4.35,
        11.2,
        0.75,
        "当前最合理的技术方向：Success-conditioned Residual Proposal for MoDA-only Planning。",
        21,
        True,
        "white",
    )
    slides.append(s)

    return prs


def make_script():
    return """# LeWM + MoDA 技术路径汇报讲稿

## 1. 开场
我这次汇报主要讲 LeWM 和 MoDA 的结合路径。LeWM 本身是 latent world model 加 CEM planning 的控制框架，MoDA 的优势是通过 depth-aware attention 做跨层信息检索。我做的工作不是简单把一个模块塞进去，而是从结构接入、训练、planning candidate pool、失败诊断，一直推进到 action proposal correction。

## 2. 任务边界
LeWM 官方或相关 checkpoint 不是只有 PushT，还包括 Cube、TwoRooms、Reacher 等任务。但我当前跑得最深的是 PushT。PushT 有连续动作、接触动力学和明显 near-miss failure，所以适合作为 MoDA-for-planning 的诊断环境。这里不能说我们已经在所有任务上验证了 MoDA，只能说框架有多任务基础，当前系统诊断主要在 PushT。

## 3. LeWM 和 MoDA
LeWM 的核心流程是：observation 经过 encoder 变成 latent state，predictor 在 latent 空间 rollout 未来状态，CEM planner 采样 action sequence 并通过 cost 选出最终 action。MoDA 的理论优势是 mixture-of-depth attention，也就是当前层可以通过 depth cache 检索前面层的 K/V 信息。这个机制放到 LeWM 里，目标是增强 predictor 的 state rollout 表达能力。

## 4. 代码接入
我主要实现了 `moda_module.py` 和 `moda_module_exact.py`。其中 exact 版本尽量保留 LeWM 原来的 conditional block 和 AdaLN-zero conditioning，只把 attention kernel 换成 MoDA attention。`MoDAAttentionExact` 在没有 cache 时退回普通 causal attention，在有 cached K/V 时调用 MoDA kernel。`MoDATransformerExact` 负责逐层收集 K/V，构造 depth cache。

训练入口包括 `train_encoder_moda.py`、`train_moda.py`、`train_moda_exact.py`。这些脚本把 MoDA predictor/encoder 接进 LeWM 的训练 pipeline，产出后续 planning evaluation 使用的 checkpoint。

## 5. Planning 诊断
接入后我们进入 planning evaluation，生成 stateroll/MoDA candidate pool。核心发现是：MoDA candidate pool 里不是没有成功候选，oracle 不低；但 raw MoDA cost 不能稳定把成功候选选到 top1。也就是说，问题不只是 coverage，而是 cost-success alignment。

一开始我尝试 BCE、calibration、selector、reranker。global AUC 看起来能涨，但 `moda_only_intra_episode_audit.py` 证明这很多是 episode difficulty leakage。episode-only 特征 AUC 都能到 0.692，但 intra-episode AUC 是 0.5；candidate-only intra-episode AUC 也只有 0.559。所以模型学到的是这个 episode 难不难，而不是同一个 episode 里哪个 candidate 应该排第一。

## 6. Baseline-safe integration
我也做过 PAC-MoDA / risk-controlled integration。它的逻辑是 baseline planner 保底，MoDA 只在 gate 判断为低风险、有机会修 baseline failure 时介入。这条线可以把系统 top1 从 81 左右提升到 82/83，说明 MoDA 有辅助候选价值。但它依赖 baseline fallback，所以不能当 MoDA-only standalone planner 结果。

## 7. Residual Proposal
当前最重要的 MoDA-only 方向是 success-conditioned residual proposal。它不再做 final rerank，而是直接修正 action proposal。具体来说，对于 raw rank0 失败但同 episode 有成功候选的样本，学习 `delta_a = success_action - raw_rank0_action`。评估时对 raw top candidates 生成 shifted candidates，再用 raw/calibrated cost 重新评估选择。

代码上，`train_data` 负责构造 residual target，`feature_np` 和 `online_feature` 负责构造离线/在线一致的特征，`fit_ridge` 学 residual，`score_candidates` 重新 rollout 打分，`run_split` 做完整评估。`moda_only_residual_confirm50_audit.py` 用固定 indices 做 paired audit，避免把不同 eval set 的 absolute top1 误读成方法变化。

## 8. 结果和口径
medium20 曾经有 raw 60 到 residual 65 的结果，但 confirm50 后更稳的结论是 raw 46 到 residual 49/50，near-miss 从 14 降到 10/11。所以不能说 MoDA-only 已经稳定 65+，但可以说 residual proposal 是目前唯一真正 MoDA-only 的正向方向：它有 consistent paired gain，并且能减少 near-miss failure。

## 9. 最终结论
一句话总结：MoDA 的候选覆盖和 depth/state rollout 有价值，但 raw cost 与 planning success 不完全对齐。我们系统排除了只靠 AUC calibration 和 final rerank 的路线，最后发现更合理的方向是从 action proposal generation 层面做 success-conditioned residual correction。
"""


if __name__ == "__main__":
    prs = make_deck()
    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)
    OUT_SCRIPT.write_text(make_script(), encoding="utf-8")
    print(OUT_PPT)
    print(OUT_SCRIPT)
